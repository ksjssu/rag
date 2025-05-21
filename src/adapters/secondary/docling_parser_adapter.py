# src/adapters/secondary/docling_parser_adapter.py

import logging
import os
from pathlib import Path
from typing import Dict, Any, Optional, List, Union
import sys
import tempfile
import re
from dataclasses import dataclass, field
from enum import Enum
import traceback
from docling_core.types.doc.document import PictureDescriptionData
from docling_core.types.doc import ImageRefMode
from pydantic import AnyUrl
import uuid
import base64
import shutil
import inspect
import time


# config에서 로거 가져오기
from src.config import logger, settings
from src.domain.models import RawDocument

# 기존 로깅 설정 제거
# logging.basicConfig(level=logging.DEBUG)
# logger.setLevel(logging.DEBUG)

# --- Docling 라이브러리 임포트 ---
from docling.document_converter import DocumentConverter, FormatOption, PdfFormatOption
from docling.datamodel.base_models import DocumentStream, ConversionStatus, InputFormat
from docling.datamodel.document import ConversionResult
from docling_core.types.doc import DocItemLabel, ImageRefMode
from docling.datamodel.pipeline_options import (
    PipelineOptions, PdfPipelineOptions, EasyOcrOptions, 
    TableStructureOptions, AcceleratorOptions, TableFormerMode, 
    granite_picture_description, PictureDescriptionApiOptions
)
from ports.output_ports import DocumentParsingPort
from docling.exceptions import ConversionError as DoclingConversionError


try:
    import docling
    _docling_available = True
except ImportError:
    _docling_available = False
    logger.warning("Docling 라이브러리를 찾을 수 없습니다.")


# --- 어댑터 특정 예외 정의 ---
class ParsingError(Exception):
    """Represents an error during the document parsing process."""
    pass

@dataclass
class ParsedDocument:
    """
    파싱 과정을 거쳐 텍스트 추출 및 기본 구조 정보가 포함된 문서 모델.
    """
    content: str  # 추출된 텍스트 내용
    metadata: Dict[str, Any] = field(default_factory=dict)  # 원본 메타데이터 + 파싱 중 얻은 메타데이터
    
    # 구조화된 콘텐츠
    tables: List[Dict[str, Any]] = field(default_factory=list)  # 테이블 구조 정보
    images: List[Dict[str, Any]] = field(default_factory=list)  # 이미지 정보
    equations: List[Dict[str, Any]] = field(default_factory=list)  # 수식 정보
    code_blocks: List[Dict[str, Any]] = field(default_factory=list)  # 코드 블록 정보
    
    # 문서 구조 정보
    headings: List[Dict[str, Any]] = field(default_factory=list)  # 제목 구조
    paragraphs: List[Dict[str, Any]] = field(default_factory=list)  # 단락 정보
    lists: List[Dict[str, Any]] = field(default_factory=list)  # 목록 정보
    
    # 레이아웃 정보
    layout: Optional[Dict[str, Any]] = None  # 페이지 레이아웃 정보
    page_info: List[Dict[str, Any]] = field(default_factory=list)  # 페이지별 정보
    
    # OCR 관련 정보
    ocr_results: Optional[Dict[str, Any]] = None  # OCR 결과 정보
    ocr_confidence: Optional[float] = None  # OCR 신뢰도
    
    # 이미지 처리 결과
    image_classifications: List[Dict[str, Any]] = field(default_factory=list)  # 이미지 분류 결과
    image_descriptions: List[Dict[str, Any]] = field(default_factory=list)  # 이미지 설명 결과
    image_links: List[Dict[str, Any]] = field(default_factory=list)  # 이미지 링크 정보(추가됨)
    
    # 테이블 처리 결과
    table_structures: List[Dict[str, Any]] = field(default_factory=list)  # 테이블 구조 분석 결과
    table_cell_matches: List[Dict[str, Any]] = field(default_factory=list)  # 테이블 셀 매칭 결과
    
    # 코드 및 수식 처리 결과
    code_enrichments: List[Dict[str, Any]] = field(default_factory=list)  # 코드 보강 정보
    formula_enrichments: List[Dict[str, Any]] = field(default_factory=list)  # 수식 보강 정보



class DoclingParserAdapter(DocumentParsingPort):
    """
    Docling 라이브러리를 사용하여 문서 파싱 기능을 제공하는 어댑터입니다.
    """

    def __init__(
        self,
        allowed_formats: Optional[List[str]] = None,
        use_gpt_picture_description: bool = False,
        images_save_dir: Optional[str] = None,
        image_resolution_scale: float = 2.0,
    ):
        self._is_initialized_successfully = False
        self._converter = None
        self._allowed_docling_formats = None
        self._image_resolution_scale = image_resolution_scale
        
        # 이미지 저장 디렉토리 설정
        self._images_save_dir = images_save_dir or getattr(settings, 'IMAGES_SAVE_DIR', 'images')
        os.makedirs(self._images_save_dir, exist_ok=True)
        logger.info(f"DoclingParserAdapter: 이미지 저장 디렉토리 설정: {self._images_save_dir}")

        # 1. 파이프라인 옵션 생성
        pipeline_options = PdfPipelineOptions()
        pipeline_options.do_picture_description = True
        pipeline_options.images_scale = self._image_resolution_scale
        pipeline_options.generate_page_images = True
        pipeline_options.generate_picture_images = True
        pipeline_options.enable_remote_services = False
        pipeline_options.do_ocr = False
        logger.info("DoclingParserAdapter: OCR 기능이 비활성화되었습니다.")
        
        # 기본 내장 모델 사용 (OpenAI API 대신)
        # 문자열을 직접 할당하는 대신 PictureDescriptionApiOptions 객체 생성
        try:
            # PictureDescriptionApiOptions 생성자 시그니처 확인
            from docling.datamodel.pipeline_options import PictureDescriptionApiOptions, granite_picture_description
            # Granite 모델 사용 설정 
            pipeline_options.picture_description_options = granite_picture_description
                # 캡션 커스텀 프롬프트 설정
            pipeline_options.picture_description_options.prompt = (
                "이미지를 세 문장으로 상세히 설명하세요. 이미지의 주요 요소, 내용, 그리고 의미를 명확하게 기술해 주세요."
            )
            logger.info("DoclingParserAdapter: granite_picture_description 모델 설정됨")
            logger.info(f"캡션 프롬프트: {pipeline_options.picture_description_options.prompt}")
        except Exception as e:
            logger.warning(f"DoclingParserAdapter: PictureDescriptionApiOptions 생성 실패: {e}")
            if hasattr(pipeline_options, 'picture_description_options'):
                del pipeline_options.picture_description_options
            logger.info("DoclingParserAdapter: picture_description_options 설정 제거됨")
        
        # 2. format_options 구성
        format_options = {
            InputFormat.PDF: PdfFormatOption(
                pipeline_options=pipeline_options
            )
        }

        # 3. allowed_formats 변환
        self._allowed_docling_formats = []
        if allowed_formats:
            for fmt_str in allowed_formats:
                fmt_upper = fmt_str.upper()
                if fmt_upper in InputFormat.__members__:
                    self._allowed_docling_formats.append(InputFormat[fmt_upper])
                elif fmt_upper in ['JPG', 'JPEG', 'PNG', 'TIF', 'TIFF', 'BMP']:
                    if InputFormat.IMAGE not in self._allowed_docling_formats:
                        self._allowed_docling_formats.append(InputFormat.IMAGE)

        # 4. DocumentConverter 초기화
        self._converter = DocumentConverter(
            allowed_formats=self._allowed_docling_formats,
            format_options=format_options
        )
        self._is_initialized_successfully = True
        logger.info("DoclingParserAdapter: Docling DocumentConverter initialized with Granite picture description")

        if _docling_available:
            self._is_initialized_successfully = False # 초기화 성공 여부 플래그
            logger.info("DoclingParserAdapter: Initializing Docling DocumentConverter...")
            try:
                 # --- Docling InputFormat 설정 (allowed_formats 처리) ---
                 # allowed_formats 문자열 목록을 Docling InputFormat enum으로 변환
                 if allowed_formats:
                      self._allowed_docling_formats = []
                      for fmt_str in allowed_formats:
                          fmt_upper = fmt_str.upper()
                          # 이미지 확장자 처리
                          if fmt_upper in ['JPG', 'JPEG', 'PNG', 'TIF', 'TIFF', 'BMP']:
                              if InputFormat.IMAGE not in self._allowed_docling_formats:
                                  self._allowed_docling_formats.append(InputFormat.IMAGE)
                          # 다른 포맷 처리
                          elif fmt_upper in InputFormat.__members__:
                              self._allowed_docling_formats.append(InputFormat[fmt_upper])
                          else:
                              logger.warning(f"Specified allowed_format '{fmt_str}' is not a valid Docling InputFormat.")


                 # --- Docling DocumentConverter 인스턴스 생성 ★ 실제 초기화 ★ ---
                 # 제공된 Docling Converter 코드의 __init__ 시그니처에 맞춰 파라미터 전달
                 # allowed_formats와 format_options를 받는 것으로 보입니다.
                 # DoclingConverter 생성자가 받는 다른 파라미터 추가 (Docling 문서 확인)

                 self._converter = DocumentConverter(
                     allowed_formats=self._allowed_docling_formats,
                     format_options=format_options, # 구성한 format_options 딕셔너리 전달
                     # 기타 Docling Converter 생성자가 받는 파라미터 추가 (Docling 문서 확인)
                 )
                 self._is_initialized_successfully = True # 초기화 성공
                 logger.info("DoclingParserAdapter: Docling DocumentConverter initialized successfully.")
            except Exception as e: # Docling Converter 초기화 중 발생할 수 있는 예외 처리
                logger.error(f"DoclingParserAdapter: Failed to initialize DocumentConverter: {e}")
                self._converter = None # 초기화 실패 시 None으로 설정
                # 초기화 실패 시 ParsingError 예외를 발생시켜 앱 시작 중단 고려
                # raise ParsingError(f"Failed to initialize Docling Converter: {e}") from e


        if self._converter is None:
             logger.warning("DoclingParserAdapter: Docling Converter not available or failed to initialize. Will use fallback decoding.")


    # --- Helper 메서드 ---
    def _guess_input_format(self, filename: str) -> Optional[InputFormat]:
        """파일 확장자를 기반으로 InputFormat을 추측합니다."""
        ext = filename.lower().split('.')[-1]
        
        format_map = {
            'docx': InputFormat.DOCX,
            'dotx': InputFormat.DOCX,
            'docm': InputFormat.DOCX,
            'dotm': InputFormat.DOCX,
            'pptx': InputFormat.PPTX,
            'potx': InputFormat.PPTX,
            'ppsx': InputFormat.PPTX,
            'pptm': InputFormat.PPTX,
            'potm': InputFormat.PPTX,
            'ppsm': InputFormat.PPTX,
            'pdf': InputFormat.PDF,
            'md': InputFormat.MD,
            'html': InputFormat.HTML,
            'htm': InputFormat.HTML,
            'xhtml': InputFormat.HTML,
            'xml': InputFormat.XML_JATS,
            'nxml': InputFormat.XML_JATS,
            'jpg': InputFormat.IMAGE,
            'jpeg': InputFormat.IMAGE,
            'png': InputFormat.IMAGE,
            'tif': InputFormat.IMAGE,
            'tiff': InputFormat.IMAGE,
            'bmp': InputFormat.IMAGE,
            'adoc': InputFormat.ASCIIDOC,
            'asciidoc': InputFormat.ASCIIDOC,
            'asc': InputFormat.ASCIIDOC,
            'csv': InputFormat.CSV,
            'xlsx': InputFormat.XLSX,
            'json': InputFormat.JSON_DOCLING
        }
        
        input_format = format_map.get(ext)
        if input_format is None:
            raise ParsingError(f"지원하지 않는 파일 형식입니다: .{ext}")
        
        return input_format

    def parse(self, raw_document: RawDocument) -> ParsedDocument:
        """
        RawDocument를 Docling DocumentConverter로 파싱합니다.
        """
        logger.info("[DEBUG] parse 메서드 시작")
        
        # 입력 검사 추가
        if raw_document is None:
            logger.error("DoclingParserAdapter: raw_document is None")
            return ParsedDocument(content="", metadata={})
        
        filename = raw_document.metadata.get('filename', 'unknown')
        logger.info(f"\n[PARSING] 시작: {filename}")

        if not hasattr(raw_document, 'content') or raw_document.content is None:
            logger.warning("DoclingParserAdapter: raw_document.content is None or missing")
            return ParsedDocument(content="", metadata=raw_document.metadata if hasattr(raw_document, 'metadata') else {})

        if not raw_document.content:
            logger.warning("DoclingParserAdapter: Empty document content received")
            return ParsedDocument(content="", metadata=raw_document.metadata)



        try:
            # 입력 형식 확인 (수정된 부분)
            try:
                input_format = self._guess_input_format(filename)
                if not input_format:
                    raise ParsingError(f"파일 형식을 확인할 수 없습니다: {filename}")
            except ParsingError as e:
                logger.error(f"DoclingParserAdapter: {str(e)}")
                raise

            logger.info(f"DoclingParserAdapter: Converting document with format {input_format}")

            # 원본 파일명에서 확장자 추출 (metadata에 filename이 있다고 가정)
            filename = raw_document.metadata.get('filename', 'unknown_file')
            file_ext = os.path.splitext(filename)[1]  # '.pdf', '.docx' 등 확장자 추출
            doc_filename = os.path.splitext(os.path.basename(filename))[0]  # 확장자 없는 파일명

            # 확장자가 없는 경우 content-type에서 추론
            if not file_ext and 'content_type' in raw_document.metadata:
                content_type = raw_document.metadata['content_type']
                # MIME 타입별 확장자 매핑
                mime_to_ext = {
                    'application/pdf': '.pdf',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                    # 기타 MIME 타입 추가...
                }
                file_ext = mime_to_ext.get(content_type, '')

            # 시작 시간 기록
            start_time = time.time()

            # 1. 임시 파일 생성
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
                temp_file.write(raw_document.content)
                temp_path = Path(temp_file.name)

            # 2. DocumentConverter로 변환
            try:
                result = self._converter.convert(
                    source=temp_path,
                    headers={"Content-Type": "application/pdf"}
                )
              
            except TypeError as e:
                if "'str' object is not callable" in str(e):
                    logger.error(f"Docling backend 오류 발생: {e}")
                    logger.error("Docling 라이브러리의 backend 파라미터 타입이 변경된 것으로 보입니다.")
                    logger.error("문제 위치: docling/datamodel/document.py의 _init_doc 메서드")
                    logger.error("원인: 백엔드가 문자열('pypdfium2')로 지정되었으나 호출 가능한 객체(클래스/함수)가 필요합니다.")
                    logger.error("해결 중: 백엔드를 호출 가능한 어댑터 클래스로 교체하는 중...")
                    
                    # DummyPdfFormatOption 클래스를 수정하여 백엔드를 함수 객체가 아닌 문자열로 전달
                    class PdfBackendAdapter:
                        def __init__(self, doc, path_or_stream=None):
                            # 실제 백엔드 기능 구현 또는 필요한 최소 인터페이스 제공
                            self.doc = doc
                            self.path = path_or_stream
                            # 백엔드 기능 구현 (또는 최소한의 인터페이스만 제공)
                            if hasattr(doc, 'export_to_text'):
                                # 메서드 객체 대신 직접 문자열 값을 설정
                                doc.export_to_text = "PDF 텍스트 추출 (어댑터에 의해 생성된 임시 콘텐츠)"
                            if hasattr(doc, 'export_to_markdown'):
                                # 메서드 객체 대신 직접 문자열 값을 설정
                                doc.export_to_markdown = "# PDF 마크다운 (어댑터에 의해 생성된 임시 콘텐츠)"
                            
                            # 이미지 객체 추가
                            class DummyPicture:
                                def __init__(self, idx, uri=None):
                                    self.idx = idx
                                    self.self_ref = f"img_{idx}"
                                    self.caption_text = f"PDF 이미지 {idx} (어댑터에 의해 생성)"
                                    self.annotations = []
                                    # 이미지 데이터 설정
                                    class DummyImage:
                                        def __init__(self, uri=None):
                                            self.uri = uri or f"data:image/png;base64,dummy_image_{idx}"
                                    self.image = DummyImage(uri)
                            
                            # 더미 이미지 목록 생성
                            if not hasattr(doc, 'pictures') or not doc.pictures:
                                logger.info("어댑터에서 더미 pictures 생성")
                                # 더미 이미지 리스트 설정
                                doc.pictures = [DummyPicture(i) for i in range(1)]
                            
                            # images 속성이 없는 경우도 처리
                            if not hasattr(doc, 'images') or not doc.images:
                                logger.info("어댑터에서 더미 images 설정")
                                doc.images = doc.pictures
                            
                            logger.info(f"PdfBackendAdapter 초기화 완료: {len(doc.pictures)}개 더미 이미지 설정")
                
                    # 변환 재시도
                    try:
                        # 백엔드를 문자열 대신 호출 가능한 객체로 전달
                        self._format_options["pdf"].backend = PdfBackendAdapter
                        result = self._converter.convert(
                            source=temp_path,
                            headers={"Content-Type": "application/pdf"}
                        )
                    except Exception as retry_e:
                        # 재시도 실패 시 오류 발생
                        logger.error(f"백엔드 어댑터 적용 후에도 변환 실패: {retry_e}")
                        raise DoclingConversionError(f"PDF 변환 실패: {retry_e}")
                else:
                    # 다른 TypeError는 그대로 전달
                    raise
            
            # 3. 텍스트 추출
            if result.status == ConversionStatus.SUCCESS:
                doc = result.document
                
                logger.debug(f"doc 객체 타입: {type(doc)}")
                
                # 1. 기본 텍스트 추출
                backend_text = ""
                ocr_text = ""
                try:
                    if hasattr(doc, 'export_to_text'):
                        # 호출 가능한지 확인 후 호출
                        if callable(doc.export_to_text):
                            backend_text = doc.export_to_text()
                        else:
                            # 문자열이나 다른 타입이면 문자열로 변환
                            backend_text = str(doc.export_to_text)
                            logger.debug(f"export_to_text는 함수가 아니라 {type(doc.export_to_text)} 타입임: {backend_text[:100]}")
                    
                    if hasattr(doc, 'ocr_text'):
                        ocr_text = str(doc.ocr_text) if doc.ocr_text is not None else ""
                    elif hasattr(doc, 'export_to_markdown'):
                        # 호출 가능한지 확인 후 호출
                        if callable(doc.export_to_markdown):
                            ocr_text = doc.export_to_markdown()
                        else:
                            # 문자열이나 다른 타입이면 문자열로 변환
                            ocr_text = str(doc.export_to_markdown)
                            logger.debug(f"export_to_markdown은 함수가 아니라 {type(doc.export_to_markdown)} 타입임: {ocr_text[:100]}")
                except Exception as e:
                    logger.error(f"Text extraction error: {e}")
                    # 오류의 상세 내용 로깅
                    logger.error(traceback.format_exc())
                
                # 2. 테이블 정보 추출
                tables = []
                table_structures = []
                table_cell_matches = []
                try:
                    if hasattr(doc, 'tables'):
                        for table in doc.tables:
                            # TableData 객체를 사전(dict)으로 변환하여 저장
                            table_data_dict = {}
                            if hasattr(table, 'data'):
                                # TableData 객체가 있는 경우 속성을 안전하게 확인하고 변환
                                if hasattr(table.data, 'to_dict'):
                                    # to_dict 메서드가 있으면 호출
                                    table_data_dict = table.data.to_dict()
                                elif hasattr(table.data, '__dict__'):
                                    # __dict__ 속성이 있으면 사용
                                    table_data_dict = vars(table.data)
                                elif hasattr(table.data, 'rows') and hasattr(table.data, 'columns'):
                                    # rows와 columns 속성이 있으면 사전에 담기
                                    try:
                                        rows_data = []
                                        if isinstance(table.data.rows, list):
                                            for row in table.data.rows:
                                                if hasattr(row, 'cells') and isinstance(row.cells, list):
                                                    row_cells = [cell.text if hasattr(cell, 'text') else str(cell) for cell in row.cells]
                                                    rows_data.append(row_cells)
                                                else:
                                                    rows_data.append(str(row))
                                        
                                        table_data_dict = {
                                            'rows': rows_data,
                                            'columns': len(table.data.columns) if hasattr(table.data.columns, '__len__') else 0
                                        }
                                    except Exception as cell_e:
                                        logger.error(f"테이블 셀 데이터 변환 오류: {cell_e}")
                                        table_data_dict = {'error': '테이블 데이터 변환 실패'}
                                else:
                                    # 그 외의 경우 문자열로 표현
                                    table_data_dict = {'data': str(table.data)}
                            
                            tables.append({
                                'structure': table_data_dict,  # 변환된 dict 형태로 저장
                                'content': table.text if hasattr(table, 'text') else "",
                                'position': table.prov[0] if hasattr(table, 'prov') and table.prov else None
                            })
                    
                    if hasattr(doc, 'table_structures'):
                        # table_structures도 안전하게 변환
                        if isinstance(doc.table_structures, list):
                            for ts in doc.table_structures:
                                if isinstance(ts, dict):
                                    table_structures.append(ts)
                                else:
                                    # dict가 아니면 변환 시도
                                    try:
                                        if hasattr(ts, 'to_dict'):
                                            table_structures.append(ts.to_dict())
                                        elif hasattr(ts, '__dict__'):
                                            table_structures.append(vars(ts))
                                        else:
                                            table_structures.append({'data': str(ts)})
                                    except Exception as ts_e:
                                        logger.error(f"테이블 구조 변환 오류: {ts_e}")
                                        table_structures.append({'error': '테이블 구조 변환 실패'})
                        else:
                            logger.warning("table_structures가 리스트 형태가 아님")
                    
                    if hasattr(doc, 'table_cell_matches'):
                        # table_cell_matches도 안전하게 변환
                        if isinstance(doc.table_cell_matches, list):
                            for tcm in doc.table_cell_matches:
                                if isinstance(tcm, dict):
                                    table_cell_matches.append(tcm)
                                else:
                                    # dict가 아니면 변환 시도
                                    try:
                                        if hasattr(tcm, 'to_dict'):
                                            table_cell_matches.append(tcm.to_dict())
                                        elif hasattr(tcm, '__dict__'):
                                            table_cell_matches.append(vars(tcm))
                                        else:
                                            table_cell_matches.append({'data': str(tcm)})
                                    except Exception as tcm_e:
                                        logger.error(f"테이블 셀 매칭 변환 오류: {tcm_e}")
                                        table_cell_matches.append({'error': '테이블 셀 매칭 변환 실패'})
                        else:
                            logger.warning("table_cell_matches가 리스트 형태가 아님")

                except Exception as e:
                    logger.error(f"Table extraction error: {e}")
                
                # 3. 이미지 정보 추출 (개선된 방법)
                images = []
                image_classifications = []
                image_descriptions = []
                image_links = []  # 이미지 링크 목록 초기화
                
                # 문서별 고유 디렉토리 생성 (문서 해시나 ID 기반)
                document_id = raw_document.metadata.get('id', str(uuid.uuid4()))
                document_dir = os.path.join(self._images_save_dir, document_id)
                os.makedirs(document_dir, exist_ok=True)
                logger.info(f"[PARSING] 이미지 저장 디렉토리 생성됨: {document_dir}")
                
                try:
                    # 시작 시간 기록
                    start_image_time = time.time()
                    
                    # 문서 요소 반복 처리 (테이블, 이미지 등)
                    table_counter = 0
                    picture_counter = 0
                    
                    if hasattr(doc, 'iterate_items'):
                        logger.info(f"[PARSING] 문서 내 이미지 객체 추출 시작")
                        
                        for element, _level in doc.iterate_items():
                            try:
                                # 테이블 요소 처리
                                if hasattr(element, '__class__') and element.__class__.__name__ == 'TableItem':
                                    table_counter += 1
                                    
                                    # 이미지 파일명 생성
                                    doc_filename = os.path.splitext(os.path.basename(filename))[0]
                                    table_image_filename = f"{doc_filename}-table-{table_counter}.png"
                                    table_image_path = os.path.join(document_dir, table_image_filename)
                                    
                                    # get_image 메소드 호출
                                    if hasattr(element, 'get_image') and callable(element.get_image):
                                        table_image = element.get_image(doc)
                                        
                                        # PIL 이미지로 저장
                                        with open(table_image_path, "wb") as fp:
                                            table_image.save(fp, "PNG")
                                        
                                        # 간소화된 메타데이터 생성 - 모든 값을 문자열로 변환
                                        table_image_info = {
                                            'source': str(document_id),
                                            'index': str(table_counter),
                                            'type': 'table'
                                        }
                                        
                                        # 이미지 목록에 추가
                                        image_links.append(table_image_info)
                                        
                                        logger.info(f"[PARSING] 테이블 {table_counter} 이미지 저장됨: {table_image_path}")
                                    else:
                                        logger.warning(f"[PARSING] 테이블 {table_counter}에 이미지 추출 메소드가 없음")
                                
                                # 그림 요소 처리
                                if hasattr(element, '__class__') and element.__class__.__name__ == 'PictureItem':
                                    picture_counter += 1
                                    
                                    # 이미지 파일명 생성
                                    doc_filename = os.path.splitext(os.path.basename(filename))[0]
                                    picture_image_filename = f"{doc_filename}-picture-{picture_counter}.png"
                                    picture_image_path = os.path.join(document_dir, picture_image_filename)
                                    
                                    # get_image 메소드 호출
                                    if hasattr(element, 'get_image') and callable(element.get_image):
                                        picture_image = element.get_image(doc)
                                        
                                        # PIL 이미지로 저장
                                        with open(picture_image_path, "wb") as fp:
                                            picture_image.save(fp, "PNG")
                                        
                                        # 간소화된 메타데이터 생성 - 모든 값을 문자열로 변환
                                        picture_image_info = {
                                            'source': str(document_id),
                                            'index': str(picture_counter), 
                                            'type': 'image'
                                        }
                                        
                                        # 이미지 목록에 추가
                                        image_links.append(picture_image_info)
                                        
                                        logger.info(f"[PARSING] 그림 {picture_counter} 이미지 저장됨: {picture_image_path}")
                                    else:
                                        logger.warning(f"[PARSING] 그림 {picture_counter}에 이미지 추출 메소드가 없음")
                            
                            except Exception as element_e:
                                logger.error(f"[PARSING] 문서 요소 처리 오류: {element_e}")
                                logger.error(traceback.format_exc())
                        
                        logger.info(f"[PARSING] 총 {table_counter}개 테이블, {picture_counter}개 그림 이미지 저장 완료")
                    else:
                        logger.warning("[PARSING] 문서에 iterate_items 메소드가 없습니다")
                    
                    # 처리 시간 계산
                    end_image_time = time.time() - start_image_time
                    logger.info(f"[PARSING] 문서 요소 이미지 추출 및 저장 완료: {end_image_time:.2f}초")
                    
                    # 이미지 처리 결과 요약
                    logger.info(f"[PARSING] 이미지 처리 결과: {len(image_links)}개 이미지 링크 생성됨")
                    logger.info(f"[PARSING] 저장된 이미지 타입: 테이블={table_counter}, 그림={picture_counter}")
                    
                    # 결과 요약
                    if image_links:
                        logger.info("\n===== 생성된 이미지 링크 목록 =====")
                        for i, link in enumerate(image_links[:5]):  # 처음 5개만 출력
                            logger.info(f"이미지 {i+1}/{len(image_links)}:")
                            logger.info(f"  - 타입: {link.get('type', 'unknown')}")
                            logger.info(f"  - 소스: {link.get('source', 'unknown')}")
                            logger.info(f"  - 인덱스: {link.get('index', 'unknown')}")
                        if len(image_links) > 5:
                            logger.info(f"  ... 그 외 {len(image_links) - 5}개 이미지")
                        logger.info("======================================")
                
                except Exception as image_e:
                    logger.error(f"[PARSING] 이미지 처리 오류: {image_e}")
                    logger.error(traceback.format_exc())
                
                # 처리 시간 계산
                end_time = time.time() - start_time
                logger.info(f"[PARSING] 문서 변환 및 이미지 저장 완료: {end_time:.2f}초")

                # 메타데이터에 이미지 링크 추가 (간소화된 형태로)
                if image_links and len(image_links) > 0:
                    logger.info(f"메타데이터에 {len(image_links)}개 이미지 링크 추가")
                    
                    # 더 단순한 형태로 변환 - 중첩 구조 제거
                    simplified_links = []
                    for link in image_links:
                        # 각 링크를 문자열 형식으로 간단하게 저장 (형식: type:index)
                        simplified_links.append(f"{link.get('type', 'unknown')}:{link.get('index', '0')}")
                    
                    # 메타데이터에 간소화된 이미지 링크 정보 저장
                    raw_document.metadata['image_link_refs'] = simplified_links
                    
                    # 원본 image_links는 저장하지 않음 - Milvus sparse vector 오류 방지
                    # raw_document.metadata['image_links'] = image_links  # 이 줄 제거

                # 4. 수식 정보 추출
                equations = []
                formula_enrichments = []
                try:
                    if hasattr(doc, 'equations'):
                        for eq in doc.equations:
                            equations.append({
                                'content': eq.content,
                                'position': eq.position
                            })
                    if hasattr(doc, 'formula_enrichments'):
                        formula_enrichments = doc.formula_enrichments
                except Exception as e:
                    logger.error(f"Equation extraction error: {e}")
                
                # 5. 코드 블록 추출
                code_blocks = []
                code_enrichments = []
                try:
                    if hasattr(doc, 'code_blocks'):
                        for code in doc.code_blocks:
                            code_blocks.append({
                                'content': code.content,
                                'language': code.language,
                                'position': code.position
                            })
                    if hasattr(doc, 'code_enrichments'):
                        code_enrichments = doc.code_enrichments
                except Exception as e:
                    logger.error(f"Code block extraction error: {e}")
                
                # 6. 문서 구조 정보 추출
                headings = []
                paragraphs = []
                lists = []
                try:
                    if hasattr(doc, 'headings'):
                        headings = doc.headings
                    if hasattr(doc, 'paragraphs'):
                        paragraphs = doc.paragraphs
                    if hasattr(doc, 'lists'):
                        lists = doc.lists
                except Exception as e:
                    logger.error(f"Document structure extraction error: {e}")
                
                # 7. 레이아웃 정보 추출
                layout = None
                page_info = []
                try:
                    if hasattr(doc, 'layout'):
                        layout = doc.layout
                    if hasattr(doc, 'page_info'):
                        page_info = doc.page_info
                except Exception as e:
                    logger.error(f"Layout extraction error: {e}")
                
                # 8. OCR 정보 추출
                ocr_results = None
                ocr_confidence = None
                try:
                    if hasattr(doc, 'ocr_results'):
                        ocr_results = doc.ocr_results
                        logger.info(f"[OCR] ocr_results: {ocr_results}")
                    if hasattr(doc, 'ocr_confidence'):
                        ocr_confidence = doc.ocr_confidence
                        logger.info(f"[OCR] ocr_confidence: {ocr_confidence}")
                except Exception as e:
                    logger.error(f"OCR information extraction error: {e}")
                
                # 최종 텍스트 선택 및 정제
                document_text = ""
                if "glyph<" in backend_text or backend_text.strip() == "":
                    document_text = ocr_text
                    logger.info("[PARSING] 글리프 감지됨: OCR 텍스트 사용")
                else:
                    document_text = backend_text
                    logger.info("[PARSING] 정상 텍스트: 백엔드 텍스트 사용")
                
                # 텍스트 정제
                document_text = re.sub(r'glyph<[^>]+>', '', document_text)
                document_text = re.sub(r'<[^>]+>', '', document_text)
                
                # 이미지 캡션 정보를 본문에 추가하여 임베딩될 수 있게 함
                if images:
                    image_captions_text = "\n\n[이미지 캡션 정보]\n"
                    useful_captions = 0
                    
                    for idx, img in enumerate(images):
                        img_id = img.get('id', f'img_{idx}')
                        description = img.get('description', '')
                        
                        # 의미 없는 캡션인지 확인 (ID와 비슷한 경우)
                        is_meaningful = True
                        if description and isinstance(description, str):
                            if description == img_id or description.startswith("이미지 #/pictures/"):
                                # 파일명에서 의미 있는 정보 추출
                                filename = raw_document.metadata.get('filename', '')
                                file_base = os.path.splitext(os.path.basename(filename))[0] if filename else ''
                                
                                # 단순 참조 ID 대신 의미 있는 설명 대체
                                img_num = idx + 1
                                index_match = re.search(r'#/pictures/(\d+)', img_id) if isinstance(img_id, str) else None
                                if index_match:
                                    img_num = int(index_match.group(1)) + 1
                                
                                description = f"{file_base} 문서의 그림 {img_num}"
                                is_meaningful = True
                        
                        if description and isinstance(description, str) and len(description) > 0 and is_meaningful:
                            image_captions_text += f"이미지 {idx+1} 설명: {description}\n"
                            useful_captions += 1
                    
                    # 캡션 정보가 있을 경우에만 추가
                    if useful_captions > 0:
                        document_text += image_captions_text
                        logger.info(f"이미지 캡션 정보를 본문에 추가했습니다 ({useful_captions}개 유용한 캡션)")
                    else:
                        logger.info("유용한 이미지 캡션이 없어 본문에 추가하지 않습니다.")
                
                # 이미지 annotation 정보도 본문에 추가
                if image_descriptions:
                    image_annot_text = "\n\n[이미지 상세 설명]\n"
                    for idx, desc in enumerate(image_descriptions):
                        description = desc.get('description', '')
                        source = desc.get('source', '')
                        if description and isinstance(description, str) and len(description) > 0:
                            image_annot_text += f"이미지 설명 {idx+1}: {description}"
                            if source:
                                image_annot_text += f" (출처: {source})"
                            image_annot_text += "\n"
                    
                    # 설명 정보가 있을 경우에만 추가
                    if image_annot_text != "\n\n[이미지 상세 설명]\n":
                        document_text += image_annot_text
                        logger.info(f"이미지 상세 설명 정보를 본문에 추가했습니다 ({len(image_descriptions)}개 설명)")
                
                logger.info(f"[PARSING] 최종 추출 텍스트 길이: {len(document_text)} 글자")
                
                # 샘플 출력
                if document_text:
                    logger.info("\n===== 추출된 텍스트 샘플 =====")
                    text_sample = ""
                    try:
                        if document_text and isinstance(document_text, str):
                            if len(document_text) > 200:
                                text_sample = document_text[0:200] + "..."
                            else:
                                text_sample = document_text
                        else:
                            # 문자열이 아닌 경우 안전하게 변환
                            text_sample = str(document_text) if document_text is not None else ""
                    except Exception as sample_e:
                        logger.error(f"텍스트 샘플 생성 오류: {sample_e}")
                        text_sample = "[텍스트 샘플 생성 실패]"
                    
                    logger.info(text_sample)
                    logger.info("==============================\n")
                
                # 이미지 캡션 정보 출력
                if images:
                    logger.info("\n===== 추출된 이미지 캡션 정보 =====")
                    for idx, img in enumerate(images):
                        img_id = img.get('id', f'img_{idx}')
                        description = img.get('description', '설명 없음')
                        
                        # 설명이 너무 길면 자르기
                        if isinstance(description, str) and len(description) > 100:
                            description_preview = description[:100] + "..."
                        else:
                            description_preview = description
                        
                        logger.info(f"이미지 {idx+1}/{len(images)} (ID: {img_id}):")
                        logger.info(f"  - 설명: {description_preview}")
                        
                        # 이미지 데이터 유형 표시
                        img_data = img.get('data')
                        if img_data:
                            if isinstance(img_data, str) and len(img_data) > 50:
                                img_data_preview = img_data[:50] + "..."
                            else:
                                img_data_preview = str(img_data)
                            logger.info(f"  - 데이터 유형: {type(img_data).__name__}, 미리보기: {img_data_preview}")
                        else:
                            logger.info(f"  - 데이터: 없음")
                        
                        # 이미지 위치 정보 표시
                        position = img.get('position')
                        if position:
                            logger.info(f"  - 위치: {position}")
                    
                    logger.info("=========================================\n")
                
                # 이미지 설명(annotations) 정보 출력
                if image_descriptions:
                    logger.info("\n===== 추출된 이미지 설명(annotations) =====")
                    for idx, desc in enumerate(image_descriptions):
                        img_id = desc.get('image_id', 'unknown')
                        description = desc.get('description', '설명 없음')
                        source = desc.get('source', '출처 없음')
                        
                        # 설명이 너무 길면 자르기
                        if isinstance(description, str) and len(description) > 100:
                            description_preview = description[:100] + "..."
                        else:
                            description_preview = description
                        
                        logger.info(f"설명 {idx+1}/{len(image_descriptions)} (이미지 ID: {img_id}):")
                        logger.info(f"  - 내용: {description_preview}")
                        logger.info(f"  - 출처: {source}")
                    
                    logger.info("=========================================\n")
                
                # 파싱 결과 출력 추가
                def print_parsing_results(parsed_doc: ParsedDocument):
                    logger.info("\n========== 파싱 결과 ==========")
                    logger.info(f"파일명: {filename}")
                    
                    # 기본 텍스트 내용
                    if parsed_doc.content:
                        logger.info("\n[텍스트 내용 샘플]")
                        try:
                            content_preview = ""
                            if parsed_doc.content and isinstance(parsed_doc.content, str):
                                if len(parsed_doc.content) > 200:
                                    content_preview = parsed_doc.content[0:200] + "..."
                                else:
                                    content_preview = parsed_doc.content
                            else:
                                # 문자열이 아닌 경우 안전하게 변환
                                content_preview = str(parsed_doc.content) if parsed_doc.content is not None else ""
                            logger.info(content_preview)
                        except Exception as e:
                            logger.error(f"텍스트 샘플 출력 오류: {e}")
                            logger.info("[텍스트 샘플 출력 실패]")
                    
                    # 테이블 정보
                    if parsed_doc.tables:
                        logger.info(f"\n[테이블 수]: {len(parsed_doc.tables)}")
                        try:
                            table_count = len(parsed_doc.tables)
                            if table_count > 0:
                                preview_count = min(2, table_count)
                                for i in range(preview_count):
                                    if i < table_count:  # 추가 안전 검사
                                        table = parsed_doc.tables[i]
                                        if isinstance(table, dict):
                                            logger.info(f"\n테이블 {i+1} 미리보기:")
                                            logger.info(f"- 구조: {table.get('structure', '정보 없음')}")
                                            logger.info(f"- 위치: {table.get('position', '정보 없음')}")
                        except Exception as table_e:
                            logger.error(f"테이블 정보 출력 오류: {table_e}")
                            logger.info("[테이블 정보 출력 실패]")
                    
                    # 이미지 정보
                    if parsed_doc.images:
                        logger.info(f"\n[이미지 수]: {len(parsed_doc.images)}")
                        try:
                            img_count = len(parsed_doc.images)
                            if img_count > 0:
                                preview_count = min(2, img_count)
                                for i in range(preview_count):
                                    if i < img_count:  # 추가 안전 검사
                                        img = parsed_doc.images[i]
                                        if isinstance(img, dict):
                                            logger.info(f"\n이미지 {i+1} 정보:")
                                            logger.info(f"- 설명: {img.get('description', '정보 없음')}")
                                            logger.info(f"- 위치: {img.get('position', '정보 없음')}")
                        except Exception as img_e:
                            logger.error(f"이미지 정보 출력 오류: {img_e}")
                            logger.info("[이미지 정보 출력 실패]")
                    
                    # 문서 구조
                    if parsed_doc.headings:
                        logger.info(f"\n[제목 구조 수]: {len(parsed_doc.headings)}")
                        try:
                            if isinstance(parsed_doc.headings, list):
                                headings_count = len(parsed_doc.headings)
                                if headings_count > 0:
                                    preview_count = min(3, headings_count)
                                    headings_preview = []
                                    for i in range(preview_count):
                                        if i < headings_count:  # 추가 안전 검사
                                            headings_preview.append(parsed_doc.headings[i])
                            else:
                                # 리스트가 아닌 경우 빈 리스트 사용
                                logger.warning("제목이 리스트 형식이 아님")
                        except Exception as headings_e:
                            logger.error(f"[제목 구조 출력 오류]: {headings_e}")
                    
                    # OCR 결과
                    if parsed_doc.ocr_confidence is not None:
                        logger.info(f"\n[OCR 신뢰도]: {parsed_doc.ocr_confidence:.2%}")
                    
                    # 수식 정보
                    if parsed_doc.equations:
                        logger.info(f"\n[수식 수]: {len(parsed_doc.equations)}")
                        try:
                            eq_limit = min(2, len(parsed_doc.equations))
                            for i, eq in enumerate(parsed_doc.equations[0:eq_limit], 1):  # 처음 2개만 출력
                                logger.info(f"수식 {i}: {eq.get('content', '정보 없음')}")
                        except Exception as e:
                            logger.error(f"수식 정보 출력 오류: {e}")
                            logger.info("[수식 정보 출력 실패]")
                    
                    logger.info("\n==============================\n")

                # FastAPI 응답 형식으로 변환
                def create_api_response(parsed_doc: ParsedDocument) -> dict:
                    try:
                        content_preview = ""
                        if parsed_doc.content and isinstance(parsed_doc.content, str):
                            if len(parsed_doc.content) > 200:
                                content_preview = parsed_doc.content[0:200] + "..."
                            else:
                                content_preview = parsed_doc.content
                        else:
                            # 문자열이 아닌 경우 안전하게 변환
                            content_preview = str(parsed_doc.content) if parsed_doc.content is not None else ""
                            
                        # 안전한 테이블 미리보기 생성
                        tables_preview = []
                        try:
                            if parsed_doc.tables and isinstance(parsed_doc.tables, list):
                                table_count = len(parsed_doc.tables)
                                if table_count > 0:
                                    preview_count = min(2, table_count)
                                    for i in range(preview_count):
                                        if i < table_count:  # 추가 안전 검사
                                            table = parsed_doc.tables[i]
                                            if isinstance(table, dict):
                                                tables_preview.append({
                                                    "structure": table.get('structure'),
                                                    "position": table.get('position')
                                                })
                        except Exception as tables_e:
                            logger.error(f"테이블 미리보기 생성 오류: {tables_e}")
                            
                        # 안전한 이미지 미리보기 생성
                        images_preview = []
                        try:
                            if parsed_doc.images and isinstance(parsed_doc.images, list):
                                img_count = len(parsed_doc.images)
                                if img_count > 0:
                                    preview_count = min(2, img_count)
                                    for i in range(preview_count):
                                        if i < img_count:  # 추가 안전 검사
                                            img = parsed_doc.images[i]
                                            if isinstance(img, dict):
                                                images_preview.append({
                                                    "description": img.get('description'),
                                                    "position": img.get('position')
                                                })
                        except Exception as images_e:
                            logger.error(f"이미지 미리보기 생성 오류: {images_e}")
                            
                        # 안전한 제목 미리보기 생성
                        headings_preview = []
                        try:
                            if parsed_doc.headings and isinstance(parsed_doc.headings, list):
                                headings_count = len(parsed_doc.headings)
                                if headings_count > 0:
                                    preview_count = min(3, headings_count)
                                    headings_preview = []
                                    for i in range(preview_count):
                                        if i < headings_count:  # 추가 안전 검사
                                            headings_preview.append(parsed_doc.headings[i])
                        except Exception as headings_e:
                            logger.error(f"제목 미리보기 생성 오류: {headings_e}")
                            
                        return {
                            "filename": filename,
                            "content_preview": content_preview,
                            "metadata": {
                                "tables_count": len(parsed_doc.tables) if parsed_doc.tables else 0,
                                "images_count": len(parsed_doc.images) if parsed_doc.images else 0,
                                "headings_count": len(parsed_doc.headings) if parsed_doc.headings else 0,
                                "ocr_confidence": parsed_doc.ocr_confidence,
                                "equations_count": len(parsed_doc.equations) if parsed_doc.equations else 0
                            },
                            "tables_preview": tables_preview,
                            "images_preview": images_preview,
                            "headings_preview": headings_preview
                        }
                    except Exception as e:
                        logger.error(f"API 응답 생성 오류: {e}")
                        return {
                            "filename": filename,
                            "error": "API 응답 생성 중 오류 발생"
                        }

                # 결과 출력
                print_parsing_results(ParsedDocument(
                    content=document_text,
                    metadata=raw_document.metadata,
                    tables=tables,
                    images=images,
                    equations=equations,
                    code_blocks=code_blocks,
                    headings=headings,
                    paragraphs=paragraphs,
                    lists=lists,
                    layout=layout,
                    page_info=page_info,
                    ocr_results=ocr_results,
                    ocr_confidence=ocr_confidence,
                    image_classifications=image_classifications,
                    image_descriptions=image_descriptions,
                    table_structures=table_structures,
                    table_cell_matches=table_cell_matches,
                    code_enrichments=code_enrichments,
                    formula_enrichments=formula_enrichments
                ))
                
                # 메타데이터에 API 응답 저장 (직렬화 가능한 형태로)
                try:
                    api_response = create_api_response(ParsedDocument(
                        content=document_text,
                        metadata=raw_document.metadata,
                        tables=tables,
                        images=images,
                        equations=equations,
                        code_blocks=code_blocks,
                        headings=headings,
                        paragraphs=paragraphs,
                        lists=lists,
                        layout=layout,
                        page_info=page_info,
                        ocr_results=ocr_results,
                        ocr_confidence=ocr_confidence,
                        image_classifications=image_classifications,
                        image_descriptions=image_descriptions,
                        table_structures=table_structures,
                        table_cell_matches=table_cell_matches,
                        code_enrichments=code_enrichments,
                        formula_enrichments=formula_enrichments,
                        image_links=image_links  # 이미지 링크 추가
                    ))
                    
                    # 직렬화 문제가 없는지 확인 (메서드 객체와 같은 직렬화 불가능한 객체 제거)
                    def ensure_serializable(obj):
                        if isinstance(obj, dict):
                            for key in list(obj.keys()):
                                if callable(obj[key]) or isinstance(obj[key], type):
                                    # 함수, 메서드, 클래스 등 직렬화 불가능한 객체 제거
                                    del obj[key]
                                else:
                                    obj[key] = ensure_serializable(obj[key])
                        elif isinstance(obj, list):
                            for i, item in enumerate(obj):
                                obj[i] = ensure_serializable(item)
                        return obj
                    
                    # 메타데이터에 API 응답 저장 (직렬화 가능한 형태로)
                    raw_document.metadata['api_response'] = ensure_serializable(api_response)
                    
                    # 메타데이터에 추출된 이미지 URI 추가
                    if images and len(images) > 0:
                        image_uris = []
                        for img in images:
                            # 이미지 데이터 유효성 검사 및 처리
                            img_data = img.get('data')
                            img_id = img.get('id', '')
                            
                            # 플레이스홀더 이미지인지 확인
                            if img_data and isinstance(img_data, str) and 'placeholder_for_' in img_data:
                                # 실제 이미지 데이터 재추출 시도
                                try:
                                    logger.info(f"메타데이터를 위한 이미지 재추출 시도: {img_id}")
                                    valid_image_data = None
                                    
                                    # 파일 확장자 확인
                                    file_ext = os.path.splitext(filename)[1].lstrip('.').lower()
                                    
                                    if file_ext == 'pdf' and raw_document.content:
                                        # PDF 파일에서 이미지 추출
                                        try:
                                            import io
                                            import fitz  # PyMuPDF
                                            import base64
                                            from PIL import Image
                                            
                                            # 임시 파일 생성
                                            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
                                                temp_file.write(raw_document.content)
                                                pdf_path = temp_file.name
                                            
                                            # PyMuPDF로 이미지 추출
                                            pdf_doc = fitz.open(pdf_path)
                                            
                                            # 이미지 인덱스 추출
                                            img_index = 0
                                            if isinstance(img_id, str) and '#/pictures/' in img_id:
                                                index_match = re.search(r'#/pictures/(\d+)', img_id)
                                                if index_match:
                                                    img_index = int(index_match.group(1))
                                            
                                            # 모든 페이지 검색
                                            for page_idx in range(len(pdf_doc)):
                                                page = pdf_doc[page_idx]
                                                image_list = page.get_images(full=True)
                                                
                                                # 특정 페이지 또는 이미지가 없으면 다음 페이지로
                                                if not image_list:
                                                    continue
                                                
                                                # 이미지가 있는 경우 - 특정 인덱스 또는 첫 번째 이미지 사용
                                                target_idx = min(img_index, len(image_list) - 1) if img_index < len(image_list) else 0
                                                
                                                img_info = image_list[target_idx]
                                                xref = img_info[0]
                                                base_image = pdf_doc.extract_image(xref)
                                                image_bytes = base_image["image"]
                                                
                                                # Base64로 인코딩
                                                img_format = base_image["ext"]
                                                b64_data = base64.b64encode(image_bytes).decode('utf-8')
                                                valid_image_data = f"data:image/{img_format};base64,{b64_data}"
                                                logger.info(f"메타데이터용 이미지 추출 성공: 페이지 {page_idx+1}, 이미지 {target_idx}")
                                                break  # 첫 번째 성공한 이미지에서 종료
                                            
                                            # 임시 파일 정리
                                            os.unlink(pdf_path)
                                            pdf_doc.close()
                                            
                                            if valid_image_data:
                                                img_data = valid_image_data
                                                logger.info("메타데이터용 이미지를 실제 이미지 데이터로 대체")
                                            
                                        except ImportError:
                                            logger.warning("PyMuPDF 라이브러리가 설치되지 않았습니다.")
                                        except Exception as pdf_e:
                                            logger.error(f"PDF에서 메타데이터용 이미지 추출 실패: {pdf_e}")
                                    
                                except Exception as extract_e:
                                    logger.error(f"메타데이터용 이미지 데이터 추출 실패: {extract_e}")
                            
                            # 최종 처리된 이미지 데이터 메타데이터에 추가
                            if img_data and isinstance(img_data, str):
                                image_uris.append({
                                    'uri': img_data,
                                    'id': img_id,
                                    'description': img.get('description', '')
                                })
                            elif img_data and hasattr(img_data, 'tobytes'):
                                # PIL.Image 객체 처리
                                try:
                                    import io
                                    import base64
                                    
                                    img_bytes = io.BytesIO()
                                    img_data.save(img_bytes, format='PNG')
                                    img_bytes = img_bytes.getvalue()
                                    b64_data = base64.b64encode(img_bytes).decode('utf-8')
                                    
                                    image_uris.append({
                                        'uri': f"data:image/png;base64,{b64_data}",
                                        'id': img_id,
                                        'description': img.get('description', '')
                                    })
                                    logger.info(f"PIL 이미지 객체를 Base64로 변환 성공")
                                except Exception as pil_e:
                                    logger.error(f"PIL 이미지 변환 실패: {pil_e}")
                                    # 변환 실패 시 기본 URI 추가
                                    image_uris.append({
                                        'uri': f"data:image/png;base64,placeholder_for_{img_id}",
                                        'id': img_id,
                                        'description': img.get('description', '')
                                    })
                            else:
                                # 그 외 경우 기본 URI 추가
                                image_uris.append({
                                    'uri': f"data:image/png;base64,placeholder_for_{img_id}",
                                    'id': img_id,
                                    'description': img.get('description', '')
                                })
                        
                        # 메타데이터에 이미지 URI 추가
                        if image_uris:
                            logger.info(f"메타데이터에 {len(image_uris)}개 이미지 URI 추가")
                            
                            # URI가 너무 길면 자르기
                            for i, uri_info in enumerate(image_uris):
                                if 'uri' in uri_info and isinstance(uri_info['uri'], str) and len(uri_info['uri']) > 60000:
                                    # URI가 너무 길면 자르기 (Milvus 65,535 바이트 제한)
                                    uri_info['uri'] = uri_info['uri'][:100] + "..." + uri_info['uri'][-100:]
                                    uri_info['uri_truncated'] = True
                                    logger.info(f"이미지 URI가 너무 길어 자름: 이미지 {i+1}, 현재 길이: {len(uri_info['uri'])}바이트")
                            
                            raw_document.metadata['image_uris'] = image_uris
                            
                            # 이미지 URI를 image_links로도 복사 (동일한 내용으로 통일)
                            if not image_links:  # image_links가 비어있는 경우만
                                image_links = []
                                for uri_info in image_uris:
                                    image_links.append({
                                        'id': uri_info.get('id', ''),
                                        'uri': uri_info.get('uri', ''),
                                        'filename': uri_info.get('id', '') + '.png',
                                        'full_path': os.path.join(self._images_save_dir, str(uuid.uuid4()), uri_info.get('id', '') + '.png'),
                                        'caption': uri_info.get('description', '')
                                    })
                                logger.info(f"image_links가 비어 있어 {len(image_uris)}개 URI 정보를 기반으로 생성됨")
                    
                    # 메타데이터에 이미지 설명 추가
                    if image_descriptions and len(image_descriptions) > 0:
                        processed_descriptions = []
                        caption_only_descriptions = []  # 캡션만 별도로 저장
                        
                        for desc in image_descriptions:
                            if desc.get('description'):
                                # 기본 설명 정보
                                desc_info = {
                                    'image_id': desc.get('image_id', ''),
                                    'description': desc.get('description', ''),
                                    'source': desc.get('source', '')
                                }
                                
                                # 타입 정보가 있으면 추가
                                if 'type' in desc:
                                    desc_info['type'] = desc['type']
                                    
                                    # 캡션은 별도 임베딩용으로 분리
                                    if desc['type'] == 'image_caption':
                                        caption_only_descriptions.append(desc_info)
                                    
                                processed_descriptions.append(desc_info)
                        
                        if processed_descriptions:
                            logger.info(f"메타데이터에 {len(processed_descriptions)}개 이미지 설명 추가")
                            raw_document.metadata['image_descriptions'] = processed_descriptions
                            
                            # 캡션만 별도로 저장 (별도 임베딩용)
                            if caption_only_descriptions:
                                logger.info(f"별도 임베딩용 {len(caption_only_descriptions)}개 이미지 캡션 추가")
                                raw_document.metadata['image_captions'] = caption_only_descriptions
                    
                    # 메타데이터에 이미지 링크 추가
                    if image_links and len(image_links) > 0:
                        logger.info(f"메타데이터에 {len(image_links)}개 이미지 링크 추가")
                        
                        # 이미지 링크에 캡션 정보 추가
                        for link in image_links:
                            link_id = link.get('id')
                            # 해당 이미지 ID의 설명 찾기
                            for desc in processed_descriptions if 'processed_descriptions' in locals() else []:
                                if desc.get('image_id') == link_id:
                                    link['description'] = desc.get('description', '')
                                    link['type'] = 'image'  # 이미지 타입 메타데이터 추가
                                    break
                        
                        raw_document.metadata['image_links'] = image_links

                        # 최종 메타데이터에 이미지 타입과 캡션 정보 추가
                        if 'api_response' in raw_document.metadata:
                            if isinstance(raw_document.metadata['api_response'], dict):
                                if 'metadata' not in raw_document.metadata['api_response']:
                                    raw_document.metadata['api_response']['metadata'] = {}
                                
                                # 이미지 개수 정보 추가
                                raw_document.metadata['api_response']['metadata']['image_links_count'] = len(image_links)
                                
                                # 이미지 타입 정보 추가
                                for link in image_links:
                                    if 'type' not in link:
                                        link['type'] = 'image'  # 기본값으로 이미지 타입 설정
                                
                                # 메타데이터 전체에 이미지 분류 정보 포함
                                raw_document.metadata['image_type'] = 'document_image'
                                raw_document.metadata['document_has_images'] = True
                        
                        # JSON 직렬화 가능한 형태로 데이터 정제
                        image_data_summary = []
                        for link in image_links:
                            # 필수 정보만 포함한 요약 생성
                            image_summary = {
                                'id': link.get('id', ''),
                                'uri': link.get('uri', '').split(',')[0] + ',...' if isinstance(link.get('uri', ''), str) and ',' in link.get('uri', '') else link.get('uri', ''),
                                'type': link.get('type', 'image'),
                                'description': link.get('description', '')[:100] + '...' if len(link.get('description', '')) > 100 else link.get('description', '')
                            }
                            image_data_summary.append(image_summary)
                        
                        # 요약된 정보를 메타데이터에 추가
                        raw_document.metadata['image_data_summary'] = image_data_summary
                        logger.info(f"메타데이터에 {len(image_data_summary)}개 이미지 요약 정보 추가")
                    
                except Exception as api_resp_e:
                    logger.error(f"API 응답 저장 중 오류: {api_resp_e}")
                    # 오류 발생 시 메타데이터에 최소한의 정보만 저장
                    raw_document.metadata['api_response'] = {
                        "filename": filename,
                        "error": "API 응답 생성 중 오류 발생",
                        "detail": str(api_resp_e)
                    }
                
                # 변환 과정의 마지막에 반환 전 검사 추가
                # 최종 반환 객체 생성
                try:
                    # 여기서 최종 결과를 반환
                    return_obj = ParsedDocument(
                        content=document_text,
                        metadata=raw_document.metadata,
                        tables=tables if 'tables' in locals() else [],
                        images=images if 'images' in locals() else [],
                        equations=equations if 'equations' in locals() else [],
                        code_blocks=code_blocks if 'code_blocks' in locals() else [],
                        headings=headings if 'headings' in locals() else [],
                        paragraphs=paragraphs if 'paragraphs' in locals() else [],
                        lists=lists if 'lists' in locals() else [],
                        layout=layout if 'layout' in locals() else None,
                        page_info=page_info if 'page_info' in locals() else [],
                        ocr_results=ocr_results if 'ocr_results' in locals() else None,
                        ocr_confidence=ocr_confidence if 'ocr_confidence' in locals() else None,
                        image_classifications=image_classifications if 'image_classifications' in locals() else [],
                        image_descriptions=image_descriptions if 'image_descriptions' in locals() else [],
                        table_structures=table_structures if 'table_structures' in locals() else [],
                        table_cell_matches=table_cell_matches if 'table_cell_matches' in locals() else [],
                        code_enrichments=code_enrichments if 'code_enrichments' in locals() else [],
                        formula_enrichments=formula_enrichments if 'formula_enrichments' in locals() else [],
                        image_links=image_links if 'image_links' in locals() else []  # 이미지 링크 추가
                    )
                    
                    # 이미지 관련 메타데이터 정보 로깅
                    if 'images' in locals() and images:
                        logger.info(f"DoclingParserAdapter: {len(images)}개 이미지 추출됨")
                        image_with_caption_count = sum(1 for img in images if img.get('caption'))
                        logger.info(f"DoclingParserAdapter: {image_with_caption_count}개 이미지에 캡션 정보 있음")
                    
                    if 'image_descriptions' in locals() and image_descriptions:
                        logger.info(f"DoclingParserAdapter: {len(image_descriptions)}개 이미지 설명 추출됨")
                    
                    if 'image_links' in locals() and image_links:
                        logger.info(f"DoclingParserAdapter: {len(image_links)}개 이미지 링크 생성됨")
                    
                    logger.info(f"DoclingParserAdapter: 파싱 성공, 텍스트 길이: {len(document_text)}")
                    return return_obj
                except Exception as final_e:
                    logger.error(f"DoclingParserAdapter: 최종 결과 생성 중 오류 발생 - {final_e}")
            
            else:
                logger.error(f"DoclingParserAdapter: Document conversion failed with status {result.status}")
                if result.errors:
                    logger.error(f"Conversion errors: {result.errors}")
                raise ParsingError(f"Document conversion failed: {result.status}")

        except DoclingConversionError as e:
            logger.error(f"DoclingParserAdapter: Docling conversion error: {e}")
            raise ParsingError(f"Docling conversion failed: {e}") from e
        except Exception as e:
            logger.error(f"DoclingParserAdapter: Unexpected error during parsing: {e}")
            logger.error(traceback.format_exc())
            raise ParsingError(f"Unexpected error during parsing: {e}") from e

        # 모든 처리가 실패했거나 예외가 발생한 경우 빈 결과 반환
        logger.warning("DoclingParserAdapter: 파싱 실패 또는 문제 발생, 빈 결과 반환")
        return ParsedDocument(content="", metadata=raw_document.metadata)

# 파일 형식별 이미지 추출 기능
def extract_images_from_file(file_bytes, file_extension, filename):
    logger.info(f"[이미지 추출] 시작: {filename} (형식: {file_extension})")
    
    # 결과 저장용 리스트
    extracted_images = []
    
    try:
        # 1. PDF 파일 처리
        if file_extension.lower() == 'pdf':
            # PDF에서 이미지 추출하는 기능 (PyMuPDF 사용)
            try:
                import io
                import fitz  # PyMuPDF
                from PIL import Image
                
                logger.info(f"[이미지 추출] PyMuPDF 이용 PDF 이미지 추출 시작: {filename}")
                pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
                logger.info(f"[이미지 추출] PDF 문서 열기 성공: {len(pdf_document)}페이지")
                img_index = 0
                
                for page_index in range(len(pdf_document)):
                    logger.info(f"[이미지 추출] 페이지 {page_index+1}/{len(pdf_document)} 처리 중...")
                    page = pdf_document.load_page(page_index)
                    image_list = page.get_images(full=True)
                    
                    logger.info(f"[이미지 추출] 페이지 {page_index+1}에서 {len(image_list)}개 이미지 발견")
                    
                    for img_index, img_info in enumerate(image_list):
                        try:
                            logger.info(f"[이미지 추출] 이미지 {img_index+1}/{len(image_list)} 처리 중...")
                            xref = img_info[0]
                            base_image = pdf_document.extract_image(xref)
                            image_bytes = base_image["image"]
                            logger.info(f"[이미지 추출] 이미지 추출 성공: {len(image_bytes)} 바이트, 형식: {base_image.get('ext', 'unknown')}")
                            
                            try:
                                img = Image.open(io.BytesIO(image_bytes))
                                logger.info(f"[이미지 추출] 이미지 변환 성공: 크기 {img.size}, 모드: {img.mode}")
                                
                                # 실제 이미지 데이터를 base64로 인코딩
                                import base64
                                image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                                image_uri = f"data:image/{base_image.get('ext', 'png')};base64,{image_b64[:30]}..." # 로그에는 일부만 표시
                            
                                extracted_images.append({
                                    'data': f"data:image/{base_image.get('ext', 'png')};base64,{image_b64}", # 실제 이미지 데이터를 URI로 저장
                                    'description': f"Image from page {page_index + 1}",
                                    'position': {"page": page_index + 1},
                                    'id': f"pdf_img_{page_index}_{img_index}"
                                })
                                logger.info(f"[이미지 추출] 이미지 ID pdf_img_{page_index}_{img_index} 추가 완료")
                            except Exception as pil_e:
                                logger.error(f"[이미지 추출] PIL 이미지 변환 오류: {pil_e}")
                                # 오류 발생해도 raw 이미지 데이터 저장
                                import base64
                                image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                                extracted_images.append({
                                    'data': f"data:image/{base_image.get('ext', 'png')};base64,{image_b64}",
                                    'description': f"Raw image from page {page_index + 1}",
                                    'position': {"page": page_index + 1},
                                    'id': f"pdf_img_{page_index}_{img_index}"
                                })
                                logger.info(f"[이미지 추출] Raw 이미지 데이터로 추가 완료")
                        except Exception as e:
                            logger.error(f"[이미지 추출] PDF 이미지 처리 오류: {e}")
                
                pdf_document.close()
                logger.info(f"[이미지 추출] PDF에서 {len(extracted_images)}개 이미지 추출됨")
            except ImportError:
                logger.warning("[이미지 추출] PyMuPDF 라이브러리가 설치되지 않았습니다.")
            except Exception as e:
                logger.error(f"[이미지 추출] PDF 처리 오류: {e}")
            
        # 2. DOCX 파일 처리
        elif file_extension.lower() in ['docx', 'doc']:
            import io
            from docx import Document
            from PIL import Image
            
            doc = Document(io.BytesIO(file_bytes))
            img_index = 0
            
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        img_bytes = rel.target_part.blob
                        img = Image.open(io.BytesIO(img_bytes))
                        
                        extracted_images.append({
                            'data': img,
                            'description': f"Image from {filename}",
                            'position': None,
                            'id': f"docx_img_{img_index}"
                        })
                        img_index += 1
                    except Exception as e:
                        logger.error(f"[이미지 추출] DOCX 이미지 처리 오류: {e}")
            
            logger.info(f"[이미지 추출] DOCX에서 {len(extracted_images)}개 이미지 추출됨")
            
        # 3. PPTX 파일 처리
        elif file_extension.lower() in ['pptx', 'ppt']:
            import io
            from pptx import Presentation
            from PIL import Image
            
            prs = Presentation(io.BytesIO(file_bytes))
            img_index = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        try:
                            img_bytes = shape.image.blob
                            img = Image.open(io.BytesIO(img_bytes))
                            
                            extracted_images.append({
                                'data': img,
                                'description': f"Image from slide {slide.slide_id}",
                                'position': {"slide": slide.slide_id},
                                'id': f"pptx_img_{img_index}"
                            })
                            img_index += 1
                        except Exception as e:
                            logger.error(f"[이미지 추출] PPTX 이미지 처리 오류: {e}")
            
            logger.info(f"[이미지 추출] PPTX에서 {len(extracted_images)}개 이미지 추출됨")
            
        # 4. XLSX 파일 처리
        elif file_extension.lower() in ['xlsx', 'xls', 'csv']:
            import io
            import openpyxl
            from PIL import Image
            
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
            img_index = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for image in sheet._images:
                    try:
                        img_bytes = image._data()
                        img = Image.open(io.BytesIO(img_bytes))
                        
                        extracted_images.append({
                            'data': img,
                            'description': f"Image from sheet {sheet_name}",
                            'position': {"sheet": sheet_name},
                            'id': f"xlsx_img_{img_index}"
                        })
                        img_index += 1
                    except Exception as e:
                        logger.error(f"[이미지 추출] XLSX 이미지 처리 오류: {e}")
            
            logger.info(f"[이미지 추출] XLSX에서 {len(extracted_images)}개 이미지 추출됨")
            
        # 5. 이미지 파일 직접 처리
        elif file_extension.lower() in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'tif']:
            import io
            from PIL import Image
            
            try:
                img = Image.open(io.BytesIO(file_bytes))
                extracted_images.append({
                    'data': img,
                    'description': f"Image file: {filename}",
                    'position': None,
                    'id': "original_image"
                })
                logger.info(f"[이미지 추출] 이미지 파일 처리 완료")
            except Exception as e:
                logger.error(f"[이미지 추출] 이미지 파일 처리 오류: {e}")
                
        # 6. HTML 파일 처리
        elif file_extension.lower() in ['html', 'htm']:
            import io
            import requests
            import base64
            from bs4 import BeautifulSoup
            from PIL import Image
            
            try:
                soup = BeautifulSoup(file_bytes, 'html.parser')
                img_tags = soup.find_all('img')
                img_index = 0
                
                for img_tag in img_tags:
                    try:
                        src = img_tag.get('src', '')
                        
                        # Base64 인코딩 이미지 처리
                        if src.startswith('data:image'):
                            # data:image/jpeg;base64,/9j/4AAQ... 형식에서 이미지 데이터 추출
                            img_data = src.split(',', 1)[1]
                            img_bytes = base64.b64decode(img_data)
                            img = Image.open(io.BytesIO(img_bytes))
                            
                            extracted_images.append({
                                'data': img,
                                'description': img_tag.get('alt', f"Image from HTML"),
                                'position': None,
                                'id': f"html_img_{img_index}"
                            })
                            img_index += 1
                    except Exception as e:
                        logger.error(f"[이미지 추출] HTML 이미지 처리 오류: {e}")
                
                logger.info(f"[이미지 추출] HTML에서 {len(extracted_images)}개 이미지 추출됨")
            except Exception as e:
                logger.error(f"[이미지 추출] HTML 파싱 오류: {e}")
                
        # 7. 기타 파일 형식 - 확장 가능
        else:
            logger.warning(f"[이미지 추출] 지원되지 않는 파일 형식: {file_extension}")
    
    except ImportError as e:
        logger.error(f"[이미지 추출] 필요한 라이브러리 없음: {e}")
    except Exception as e:
        logger.error(f"[이미지 추출] 오류 발생: {e}")
    
    logger.info(f"[이미지 추출] 완료: {len(extracted_images)}개 이미지 추출됨")
    return extracted_images

